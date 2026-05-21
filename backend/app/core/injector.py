"""
Core PPTX re-injection logic.
Re-inserts translated text while preserving original styling.
"""
from pptx import Presentation
from pptx.shapes.base import BaseShape as Shape
from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.text.text import TextFrame
from pptx.util import Pt, Emu
from typing import Optional
import copy
from lxml import etree

from app.models import TranslatedRun, TranslationJob, SpatialConstraints


SMARTART_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData'
SMARTART_DGM_URI = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
PPTX_NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
}


# Average character widths for font size estimation
# These are approximations for common fonts
CHAR_WIDTH_RATIOS = {
    'ja': 1.0,   # Japanese characters (full-width)
    'en': 0.5,   # English characters (half-width)
}


def estimate_text_width(text: str, font_size: float) -> float:
    """
    Estimate text width based on character content.
    Japanese characters are typically wider than English.
    """
    # Count Japanese vs English characters
    ja_count = sum(1 for c in text if '\u3040' <= c <= '\u30ff' or '\u4e00' <= c <= '\u9fff')
    en_count = len(text) - ja_count
    
    # Estimate width
    ja_width = ja_count * font_size * CHAR_WIDTH_RATIOS['ja']
    en_width = en_count * font_size * CHAR_WIDTH_RATIOS['en']
    
    return ja_width + en_width


def check_text_fit(
    original_text: str,
    translated_text: str,
    constraints: SpatialConstraints,
    font_size: Optional[float],
) -> tuple[bool, Optional[float], Optional[str]]:
    """
    Check if translated text fits in the original text box.
    
    Returns:
        (fits, adjusted_font_size, reason)
    """
    if not font_size:
        font_size = 12.0  # Default font size
    
    box_width_emu = constraints.width
    box_width_pt = box_width_emu / 12700  # Convert EMU to points (rough approximation)
    
    original_width = estimate_text_width(original_text, font_size)
    translated_width = estimate_text_width(translated_text, font_size)
    
    # Check if translated text exceeds box width by more than 10%
    if translated_width > box_width_pt * 1.1:
        # Try reducing font size by up to 20%
        max_reduction = 0.2
        for reduction in [0.05, 0.1, 0.15, 0.2]:
            adjusted_size = font_size * (1 - reduction)
            adjusted_width = estimate_text_width(translated_text, adjusted_size)
            if adjusted_width <= box_width_pt:
                return True, adjusted_size, "overflow"
        return False, None, "overflow"
    
    # Check if text is significantly smaller (more than 50% smaller)
    if translated_width < original_width * 0.5 and translated_width < box_width_pt * 0.3:
        # Consider increasing font size (optional)
        return True, font_size, None
    
    return True, None, None


def set_run_text_safe(run, new_text: str):
    """
    Safely set text on a run, preserving all formatting.
    """
    # Store original properties
    original_font = run.font
    
    # Set the new text
    run.text = new_text
    
    # Restore font properties (they should be preserved, but let's be safe)
    try:
        if original_font.name:
            run.font.name = original_font.name
        if original_font.size:
            run.font.size = original_font.size
        if original_font.bold is not None:
            run.font.bold = original_font.bold
        if original_font.italic is not None:
            run.font.italic = original_font.italic
        if original_font.underline is not None:
            run.font.underline = original_font.underline
        if original_font.strike is not None:
            run.font.strike = original_font.strike
        if original_font.color and original_font.color.rgb:
            run.font.color.rgb = original_font.color.rgb
    except Exception:
        pass  # Some properties might not be settable


def find_shape_by_index(shape, shape_idx: str) -> Optional[Shape]:
    """
    Find a shape by its index, handling nested groups.
    shape_idx can be like "0", "1_0", "1_table_0_0"
    """
    parts = str(shape_idx).split('.')
    
    # Handle table cells
    if 'table' in shape_idx:
        # This is a table cell, handled separately
        return None
    
    try:
        # Handle grouped shapes
        if isinstance(shape, GroupShape):
            first_idx = int(parts[0])
            sub_shape = list(shape.shapes)[first_idx]
            if len(parts) > 1:
                return find_shape_by_index(sub_shape, '.'.join(parts[1:]))
            return sub_shape
    except (IndexError, ValueError):
        pass
    
    return None


def inject_smartart_text(
    shape: GraphicFrame,
    translated_runs: list[TranslatedRun],
    slide_idx: int,
) -> list[TranslatedRun]:
    """Inject translated text into a SmartArt diagram via its XML."""
    failed_runs = []
    try:
        shape_el = shape._element
        graphic_data = shape_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData')
        if graphic_data is None:
            return translated_runs
        
        uri = graphic_data.get('uri', '')
        if SMARTART_DGM_URI not in uri:
            return translated_runs
        
        # Find relId — SmartArt uses <dgm:relIds> with r:dm
        rel_ids_el = graphic_data.find('.//dgm:relIds', PPTX_NS)
        rel_id = None
        if rel_ids_el is not None:
            rel_id = rel_ids_el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm')
        else:
            for child in graphic_data:
                rel_id = child.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rel_id:
                    break
        
        if not rel_id:
            return translated_runs
        
        # Get slide part
        try:
            # Access the diagram data part via relationship
            dgm_part = shape.part.related_part(rel_id) if hasattr(shape, 'part') else None
        except (KeyError, AttributeError):
            dgm_part = None
        
        if dgm_part is None:
            return translated_runs
        
        dgm_xml = etree.fromstring(dgm_part.blob)
        
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        
        for tr in translated_runs:
            try:
                run_parts = tr.run_id.split('_')
                run_idx = int(run_parts[4]) if len(run_parts) > 4 else 0
                
                # Find the corresponding <a:t> element in the data model
                # Scan through all <dgm:pt> → <a:t> elements
                a_t_elements = []
                for pt in dgm_xml.iter(f'{{{PPTX_NS["dgm"]}}}pt'):
                    for t_elem in pt.iter(f'{{{a_ns}}}t'):
                        a_t_elements.append(t_elem)
                
                if run_idx < len(a_t_elements):
                    t_elem = a_t_elements[run_idx]
                    t_elem.text = tr.translated_text
                else:
                    failed_runs.append(tr)
            except Exception:
                failed_runs.append(tr)
        
        # Write back modified XML
        if len(failed_runs) < len(translated_runs):
            dgm_part._blob = etree.tostring(dgm_xml, xml_declaration=True, encoding='UTF-8')
        
    except Exception as e:
        print(f"[INJECTOR] SmartArt injection error: {e}")
        failed_runs.extend(translated_runs)
    
    return failed_runs


def replace_text_in_shape(
    shape: Shape,
    translated_runs: list[TranslatedRun],
    slide_idx: int,
    shape_idx: str,
) -> list[TranslatedRun]:
    """
    Replace text in a shape with translated text.
    
    Returns list of runs that couldn't be replaced.
    """
    failed_runs = []
    
    # Handle grouped shapes
    if isinstance(shape, GroupShape):
        for sub_idx, sub_shape in enumerate(shape.shapes):
            runs = replace_text_in_shape(
                sub_shape,
                translated_runs,
                slide_idx,
                f"{shape_idx}.{sub_idx}",
            )
            failed_runs.extend(runs)
        return failed_runs
    
    # Handle tables
    if isinstance(shape, GraphicFrame) and shape.has_table:
        table = shape.table
        for translated_run in translated_runs:
            try:
                run_parts = translated_run.run_id.split('_')
                shape_segments = run_parts[2].split('.') if len(run_parts) > 2 else []
                
                table_row = 0
                table_col = 0
                if len(shape_segments) >= 4 and 'table' in shape_segments:
                    table_row = int(shape_segments[2])
                    table_col = int(shape_segments[3])
                
                cell = table.rows[table_row].cells[table_col]
                paragraph_idx = int(run_parts[3]) if len(run_parts) > 3 else 0
                run_offset = int(run_parts[4]) if len(run_parts) > 4 else 0
                
                para = list(cell.text_frame.paragraphs)[paragraph_idx]
                run = list(para.runs)[run_offset]
                
                try:
                    orig_font_size = run.font.size
                except Exception:
                    orig_font_size = None
                
                set_run_text_safe(run, translated_run.translated_text)
                
                if translated_run.adjusted_font_size:
                    run.font.size = Pt(translated_run.adjusted_font_size)
                        
            except (IndexError, ValueError) as e:
                failed_runs.append(translated_run)
        return failed_runs
    
    # Handle SmartArt diagrams
    if isinstance(shape, GraphicFrame) and 'smartart' in str(translated_runs[0].run_id if translated_runs else ''):
        return inject_smartart_text(shape, translated_runs, slide_idx)
    
    # Handle regular shapes with text frames
    if not hasattr(shape, 'text_frame') or shape.text_frame is None:
        return failed_runs
    
    text_frame = shape.text_frame
    
    # Group runs by paragraph
    runs_by_paragraph = {}
    for tr in translated_runs:
        para_parts = tr.run_id.split('_')
        para_idx = para_parts[3] if len(para_parts) > 3 else '0'
        if para_idx not in runs_by_paragraph:
            runs_by_paragraph[para_idx] = []
        runs_by_paragraph[para_idx].append(tr)
    
    # Replace text in each paragraph
    paragraphs = list(text_frame.paragraphs)
    print(f"  [INJECTOR] shape {shape_idx}: {len(paragraphs)} paragraphs, {len(runs_by_paragraph)} run groups")
    
    for para_idx_str, runs in runs_by_paragraph.items():
        try:
            para_idx = int(para_idx_str)
            if para_idx >= len(paragraphs):
                continue
            
            paragraph = paragraphs[para_idx]
            paragraph_runs = list(paragraph.runs)
            
            for tr in runs:
                run_i_parts = tr.run_id.split('_')
                run_idx = int(run_i_parts[4]) if len(run_i_parts) > 4 else 0
                
                if run_idx < len(paragraph_runs):
                    run = paragraph_runs[run_idx]
                    
                    # Read original font size to preserve
                    orig_font_size = None
                    try:
                        if run.font.size:
                            orig_font_size = run.font.size
                    except Exception:
                        pass
                    
                    set_run_text_safe(run, tr.translated_text)
                    
                    # Apply font size adjustment if explicitly set (from check_text_fit logic)
                    if tr.adjusted_font_size:
                        run.font.size = Pt(tr.adjusted_font_size)
                else:
                    failed_runs.append(tr)
                    
        except (IndexError, ValueError) as e:
            failed_runs.extend(runs)
    
    return failed_runs


def inject_translations(
    input_path: str,
    output_path: str,
    translated_runs: list[TranslatedRun],
    original_document,  # PPTXDocument from extractor
) -> tuple[bool, list[TranslatedRun]]:
    """
    Inject translated text into a PPTX file.
    
    Args:
        input_path: Path to original PPTX
        output_path: Path for translated PPTX
        translated_runs: List oftranslated text runs
        original_document: Original PPTXDocument with spatial constraints
        
    Returns:
        (success, failed_runs) - list of runs that couldn't be replaced
    """
    # Open the original presentation
    prs = Presentation(input_path)
    
    failed_runs = []
    
    # Group runs by slide
    runs_by_slide = {}
    for tr in translated_runs:
        slide_idx = tr.run_id.split('_')[1] if '_' in tr.run_id else '0'
        if slide_idx not in runs_by_slide:
            runs_by_slide[slide_idx] = []
        runs_by_slide[slide_idx].append(tr)
    
    # Process each slide
    for slide_idx_str, slide_runs in runs_by_slide.items():
        try:
            slide_idx = int(slide_idx_str)
            slide = prs.slides[slide_idx]
            
            # Group runs by shape
            runs_by_shape = {}
            for tr in slide_runs:
                rid_parts = tr.run_id.split('_')
                shape_idx = rid_parts[2] if len(rid_parts) > 2 else '0'
                if shape_idx not in runs_by_shape:
                    runs_by_shape[shape_idx] = []
                runs_by_shape[shape_idx].append(tr)
            
            # Process each shape
            for shape_idx_str, shape_runs in runs_by_shape.items():
                try:
                    shape_idx_parts = shape_idx_str.split('.')
                    
                    # Handle SmartArt shapes
                    is_smartart = 'smartart' in shape_idx_str
                    
                    # Find the shape
                    shape = None
                    if is_smartart:
                        # SmartArt: find the GraphicFrame in the slide
                        for s_idx, s in enumerate(slide.shapes):
                            if isinstance(s, GraphicFrame):
                                shape = s
                                break
                    elif len(shape_idx_parts) == 1:
                        # Simple index
                        idx = int(shape_idx_parts[0])
                        shape = slide.shapes[idx]
                    else:
                        # Nested or special shape
                        first_idx = int(shape_idx_parts[0])
                        shape = slide.shapes[first_idx]
                        
                        if 'table' in shape_idx_str and isinstance(shape, GraphicFrame):
                            # Table cell - handled in replace_text_in_shape
                            pass
                        elif isinstance(shape, GroupShape):
                            # Nested group — use '.' separator to match find_shape_by_index
                            shape = find_shape_by_index(shape, '.'.join(shape_idx_parts[1:]))
                    
                    if shape:
                        failed = replace_text_in_shape(shape, shape_runs, slide_idx, shape_idx_str)
                        failed_runs.extend(failed)
                    else:
                        # Try to find by iterating all shapes
                        for s_idx, s in enumerate(slide.shapes):
                            if f"_{s_idx}_" in shape_idx_str or shape_idx_str == str(s_idx):
                                failed = replace_text_in_shape(s, shape_runs, slide_idx, str(s_idx))
                                failed_runs.extend(failed)
                                break
                                
                except Exception as e:
                    failed_runs.extend(shape_runs)
                    
        except Exception as e:
            failed_runs.extend(slide_runs)
    
    # Save the translated presentation
    print(f"[INJECTOR] Processed {len(translated_runs)} runs, {len(failed_runs)} failed")
    prs.save(output_path)
    
    if failed_runs:
        print(f"[INJECTOR] Failed runs:")
        for fr in failed_runs[:10]:
            print(f"  {fr.run_id}: {fr.original_text[:30]} -> {fr.translated_text[:30]}")
    
    return len(failed_runs) == 0, failed_runs