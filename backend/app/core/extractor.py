"""
Core PPTX extraction logic.
Extracts text runs while preserving XML styling attributes.
"""
from pptx import Presentation
from pptx.shapes.base import BaseShape as Shape
from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table
from pptx.text.text import TextFrame
from pptx.util import Pt
from lxml import etree
from typing import Optional
import uuid

from app.models import (
    FontStyle,
    SpatialConstraints,
    TextRun,
    TextBox,
    Slide,
    PPTXDocument,
)


# XML namespaces for PPTX
PPTX_NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
}

SMARTART_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData'


def get_shape_id(shape: Shape) -> int:
    """Get the unique ID of a shape."""
    return shape.shape_id


def extract_font_size(run) -> Optional[float]:
    """Extract font size from run XML, handling None case."""
    try:
        if run.font.size:
            return run.font.size.pt
        # Try to get from XML directly
        rPr = run._r.get_or_add_rPr()
        sz = rPr.get('{http://schemas.openxmlformats.org/drawingml/2006/main}sz')
        if sz:
            return int(sz) / 100  # XML stores in hundredths of a point
    except Exception:
        pass
    return None


def extract_font_color(run) -> Optional[str]:
    """Extract font color as hex string."""
    try:
        if run.font.color.rgb:
            return f"#{run.font.color.rgb}"
        # Check for theme color
        if run.font.color.theme_color:
            return f"theme:{run.font.color.theme_color}"
    except Exception:
        # Try XML extraction
        try:
            rPr = run._r.get_or_add_rPr()
            solidFill = rPr.find('.//a:solidFill', PPTX_NAMESPACES)
            if solidFill is not None:
                srgbClr = solidFill.find('a:srgbClr', PPTX_NAMESPACES)
                if srgbClr is not None:
                    return f"#{srgbClr.get('val')}"
        except Exception:
            pass
    return None


def extract_font_name(run) -> Optional[str]:
    """Extract font name/typeface."""
    try:
        if run.font.name:
            return run.font.name
        # Try XML extraction
        rPr = run._r.get_or_add_rPr()
        latin = rPr.find('a:latin', PPTX_NAMESPACES)
        if latin is not None:
            return latin.get('typeface')
    except Exception:
        pass
    return None


def extract_style_from_run(run) -> FontStyle:
    """Extract all styling attributes from a text run."""
    style = FontStyle(
        font_size=extract_font_size(run),
        font_color=extract_font_color(run),
        font_name=extract_font_name(run),
        bold=run.font.bold if run.font.bold else False,
        italic=run.font.italic if run.font.italic else False,
        underline=bool(run.font.underline) if run.font.underline else False,
        strike_through=False,  # Not available in python-pptx
    )
    
    # Check for vertical text (tategaki)
    try:
        rPr = run._r.get_or_add_rPr()
        # Vertical text is indicated by the 'vert' attribute
        if rPr.get('{http://schemas.openxmlformats.org/drawingml/2006/main}vert') == 'vert':
            style.vertical = True
    except Exception:
        pass
    
    return style


def get_spatial_constraints(shape: Shape) -> SpatialConstraints:
    """Get spatial constraints for a shape."""
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
        
        # Try to get anchor/alignment
        anchor_x = "left"
        anchor_y = "top"
        
        if hasattr(shape, 'text_frame') and shape.text_frame:
            paragraphs = list(shape.text_frame.paragraphs)
            if paragraphs:
                first_para = paragraphs[0]
                if first_para.alignment:
                    align_map = {
                        'LEFT': 'left',
                        'CENTER': 'center',
                        'RIGHT': 'right',
                    }
                    anchor_x = align_map.get(str(first_para.alignment), 'left')
        
        if hasattr(shape, 'vertical_anchor'):
            vert_map = {
                'TOP': 'top',
                'MIDDLE': 'middle',
                'BOTTOM': 'bottom',
            }
            anchor_y = vert_map.get(str(shape.vertical_anchor), 'top')
        
        return SpatialConstraints(
            left=left,
            top=top,
            width=width,
            height=height,
            anchor_x=anchor_x,
            anchor_y=anchor_y,
        )
    except Exception:
        # Return defaults
        return SpatialConstraints(left=0, top=0, width=0, height=0)


_run_counter: int = 0
_box_counter: int = 0


def generate_run_id(slide_idx: int, shape_idx: int | str, para_idx: int, run_idx: int) -> str:
    """Generate a unique ID for a text run."""
    global _run_counter
    # Use '.' to join compound shape paths so the run_id remains parseable
    shape_str = str(shape_idx).replace('_', '.')
    run_id = f"run_{slide_idx}_{shape_str}_{para_idx}_{run_idx}_{_run_counter}"
    _run_counter += 1
    return run_id


def get_xml_path(run) -> str:
    """Get the XML path for a run for later re-injection."""
    try:
        # Build a path like: /p:sld/p:sp/p:txBody/a:p/a:r
        parts = []
        element = run._r
        while element is not None:
            tag = element.tag
            # Extract local name
            if '}' in tag:
                tag = tag.split('}')[1]
            parts.append(tag)
            element = element.getparent()
        return '/'.join(reversed(parts))
    except Exception:
        return ""


def extract_runs_from_text_frame(
    text_frame: TextFrame,
    slide_idx: int,
    shape_idx: int | str,
    shape_type: str,
    run_id_shape_idx: int | str | None = None,
) -> list[TextRun]:
    """Extract all text runs from a text frame.
    
    Args:
        text_frame: The text frame to extract from
        slide_idx: Slide index
        shape_idx: Shape index stored in the model (must be int-compatible)
        shape_type: Type of shape
        run_id_shape_idx: If different from shape_idx (e.g., table cell path),
                          used for run_id generation only. Defaults to shape_idx.
    """
    runs = []
    effective_run_id_idx = run_id_shape_idx if run_id_shape_idx is not None else shape_idx
    
    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        for run_idx, run in enumerate(paragraph.runs):
            if not run.text.strip():
                continue
            
            style = extract_style_from_run(run)
            xml_path = get_xml_path(run)
            run_id = generate_run_id(slide_idx, effective_run_id_idx, para_idx, run_idx)
            
            text_run = TextRun(
                run_id=run_id,
                text=run.text,
                style=style,
                slide_index=slide_idx,
                shape_index=shape_idx,
                paragraph_index=para_idx,
                run_index=run_idx,
                xml_path=xml_path,
            )
            runs.append(text_run)
    
    return runs


# SmartArt text extraction
SMARTART_DGM_URI = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'


def extract_smartart_text(
    shape: GraphicFrame,
    slide_part,
    slide_idx: int,
    shape_idx: int | str,
) -> list[TextRun]:
    """Extract text from a SmartArt diagram shape."""
    runs = []
    try:
        shape_el = shape._element
        # Find the graphic data element — it's in the a: namespace, not p:
        graphic_data = shape_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData')
        if graphic_data is None:
            return runs
        
        # Check if this is a SmartArt (not a chart or table)
        uri = graphic_data.get('uri', '')
        if SMARTART_DGM_URI not in uri:
            return runs
        
        # Find the dgm relationship ID — SmartArt uses <dgm:relIds> with r:dm, r:lo, r:qs, r:cs
        rel_ids_el = graphic_data.find('.//dgm:relIds', PPTX_NAMESPACES)
        rel_id = None
        if rel_ids_el is not None:
            rel_id = rel_ids_el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm')
        else:
            # Fallback: try direct relId on child elements
            for child in graphic_data:
                rel_id = child.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rel_id:
                    break
        
        if not rel_id:
            return runs
        
        # Access the diagram data part via relationship
        try:
            dgm_part = slide_part.related_part(rel_id)
            dgm_xml = etree.fromstring(dgm_part.blob)
        except (KeyError, AttributeError) as e:
            return runs
        
        # Find all text elements in the diagram
        # SmartArt stores text in <dgm:t> → <a:p> → <a:r> → <a:t> elements
        # We extract from <a:t> which contains the actual text
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        text_idx = 0
        for pt in dgm_xml.iter(f'{{{PPTX_NAMESPACES["dgm"]}}}pt'):
            for t_elem in pt.iter(f'{{{a_ns}}}t'):
                text = (t_elem.text or '').strip()
                if not text:
                    continue
                
                # Build XML path for re-injection
                parts = []
                elem = t_elem
                while elem is not None:
                    tag = elem.tag
                    if '}' in tag:
                        tag = tag.split('}')[1]
                    parts.append(tag)
                    elem = elem.getparent()
                xml_path = '/'.join(reversed(parts))
                
                run_id = generate_run_id(slide_idx, f"smartart_{shape_idx}", 0, text_idx)
                text_idx += 1
                
                runs.append(TextRun(
                    run_id=run_id,
                    text=text,
                    style=FontStyle(),
                    slide_index=slide_idx,
                    shape_index=shape_idx,
                    paragraph_index=0,
                    run_index=text_idx,
                    xml_path=xml_path,
                ))
    except Exception as e:
        print(f"[EXTRACTOR] SmartArt extraction error: {e}")
    
    return runs


def extract_from_shape(
    shape: Shape,
    slide_idx: int,
    shape_idx: int | str,
    slide_part=None,  # pptx.opc.package.Part for accessing relationships
) -> list[TextBox]:
    """Extract text boxes from a shape (recursive for groups)."""
    text_boxes = []

    # Handle group shapes recursively
    if isinstance(shape, GroupShape):
        for sub_idx, sub_shape in enumerate(shape.shapes):
            nested_shape_path = f"{shape_idx}_{sub_idx}"
            text_boxes.extend(extract_from_shape(sub_shape, slide_idx, nested_shape_path, slide_part))
        return text_boxes

    # Handle graphic frames (charts, diagrams, tables)
    if isinstance(shape, GraphicFrame):
        # Check for tables first
        if shape.has_table:
            table: Table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if cell.text_frame:
                        # Use compound ID so run_id encodes row/col position
                        table_id = f"{shape_idx}_table_{row_idx}_{col_idx}"
                        runs = extract_runs_from_text_frame(
                            cell.text_frame,
                            slide_idx,
                            shape_idx,  # int for the model
                            "table_cell",
                            run_id_shape_idx=table_id,  # compound ID for run_id only
                        )
                        if runs:
                            constraints = get_spatial_constraints(shape)
                            tb_shape_idx = int(shape_idx) if isinstance(shape_idx, (int, str)) and str(shape_idx).isdigit() else 0
                            text_boxes.append(TextBox(
                                box_id=f"box_{slide_idx}_{table_id}",
                                shape_type="table_cell",
                                slide_index=slide_idx,
                                shape_index=tb_shape_idx,  # int for model
                                runs=runs,
                                constraints=constraints,
                            ))
        # SmartArt diagrams
        if slide_part is not None:
            smartart_texts = extract_smartart_text(shape, slide_part, slide_idx, shape_idx)
            if smartart_texts:
                text_boxes.append(TextBox(
                    box_id=f"box_{slide_idx}_smartart_{shape_idx}",
                    shape_type="smartart",
                    slide_index=slide_idx,
                    shape_index=shape_idx,
                    runs=smartart_texts,
                    constraints=get_spatial_constraints(shape),
                ))
        return text_boxes
    
    # Regular shapes with text frames
    if not hasattr(shape, 'text_frame'):
        return text_boxes
    
    if shape.text_frame is None:
        return text_boxes
    
    runs = extract_runs_from_text_frame(
        shape.text_frame,
        slide_idx,
        shape_idx,
        "shape",
    )
    
    if runs:
        constraints = get_spatial_constraints(shape)
        global _box_counter
        tb = TextBox(
            box_id=f"box_{slide_idx}_{_box_counter}",
            shape_type="shape",
            slide_index=slide_idx,
            shape_index=shape_idx,
            runs=runs,
            constraints=constraints,
        )
        _box_counter += 1
        text_boxes.append(tb)
    
    return text_boxes


def extract_slide(pptx: Presentation, slide_idx: int) -> Slide:
    """Extract all text from a single slide."""
    slide = pptx.slides[slide_idx]
    slide_id = slide.slide_id
    slide_part = slide.part  # For accessing relationships (SmartArt, etc.)
    
    text_boxes = []
    
    for shape_idx, shape in enumerate(slide.shapes):
        text_boxes.extend(extract_from_shape(shape, slide_idx, shape_idx, slide_part))
    
    # Calculate total runs
    total_runs = sum(len(tb.runs) for tb in text_boxes)
    
    return Slide(
        slide_index=slide_idx,
        slide_id=slide_id,
        text_boxes=text_boxes,
        preview_base64=None,  # Will be generated if requested
    )


def extract_pptx(file_path: str, generate_preview: bool = False) -> PPTXDocument:
    """
    Extract all text runs from a PPTX file.
    
    Args:
        file_path: Path to the PPTX file
        generate_preview: Whether to generate base64 preview images
        
    Returns:
        PPTXDocument with all extracted text runs
    """
    pptx = Presentation(file_path)
    
    slides = []
    total_runs = 0
    
    for slide_idx in range(len(pptx.slides)):
        slide = extract_slide(pptx, slide_idx)
        slides.append(slide)
        total_runs += sum(len(tb.runs) for tb in slide.text_boxes)
    
    filename = file_path.split('/')[-1] if '/' in file_path else file_path.split('\\')[-1]
    
    return PPTXDocument(
        filename=filename,
        slides=slides,
        total_runs=total_runs,
        extraction_metadata={
            "total_slides": len(slides),
            "total_text_boxes": sum(len(s.text_boxes) for s in slides),
            "extraction_method": "python-pptx",
        },
    )